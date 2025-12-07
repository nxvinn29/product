"""
Worker to convert PDF to PowerPoint (.pptx)
Uses pdf2image to extract pages and python-pptx to create presentation
"""
from .celery_app import celery_app
import os
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches


@celery_app.task(name="pdf_to_pptx", bind=True)
def pdf_to_pptx(self, job_id: str, input_path: str, params: dict = None) -> dict:
    """
    Convert PDF to PowerPoint presentation
    
    Args:
        job_id: Unique job identifier
        input_path: Path to input PDF
        params: dict with optional keys:
            - dpi: DPI for PDF rendering (default 150)
            - title: Presentation title
    
    Returns:
        dict with file_path to output .pptx file
    """
    if params is None:
        params = {}
    
    output_dir = os.path.dirname(input_path)
    output_path = os.path.join(output_dir, "output.pptx")
    
    try:
        dpi = params.get("dpi", 150)
        title = params.get("title", "PDF Presentation")
        
        # Convert PDF pages to images
        images = convert_from_path(input_path, dpi=dpi)
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = title
        
        # Add image slides
        for idx, image in enumerate(images, 1):
            # Use blank layout
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Save image to temp location
            temp_img_path = os.path.join(output_dir, f"page_{idx}.png")
            image.save(temp_img_path, "PNG")
            
            # Add image to slide, filling the slide
            left = Inches(0)
            top = Inches(0)
            slide.shapes.add_picture(temp_img_path, left, top, width=prs.slide_width, height=prs.slide_height)
            
            # Clean up temp image
            os.remove(temp_img_path)
        
        # Save presentation
        prs.save(output_path)
        
        return {"file_path": output_path}
    
    except Exception as e:
        raise Exception(f"PDF to PPTX conversion failed: {str(e)}")
